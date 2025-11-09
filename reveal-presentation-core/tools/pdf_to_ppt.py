#!/usr/bin/env python3
"""
将PDF的每一页转换为图片并插入到PPT中
"""

import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Emu
from io import BytesIO

def pdf_to_ppt(pdf_path, output_ppt_path):
    """
    将PDF转换为PPT

    Args:
        pdf_path: PDF文件路径
        output_ppt_path: 输出的PPT文件路径
    """
    print(f"正在读取PDF: {pdf_path}")

    # 将PDF转换为图片列表
    # dpi=200 可以根据需要调整清晰度，越高越清晰但文件越大
    # 使用 200 DPI 确保高清晰度 (1440 pts * 200/72 = 4000 像素宽)
    images = convert_from_path(pdf_path, dpi=200)

    print(f"PDF共有 {len(images)} 页")

    # 创建PPT
    prs = Presentation()

    # 设置PPT尺寸为真正的 1920x1080 高分辨率
    # PowerPoint 内部使用 EMU (English Metric Units): 1 inch = 914400 EMU
    # 1920px / 72 dpi = 26.67 inches, 1080px / 72 dpi = 15 inches
    prs.slide_width = Emu(int(1920 * 914400 / 72))   # 1920 像素
    prs.slide_height = Emu(int(1080 * 914400 / 72))  # 1080 像素

    # 遍历每一页图片
    for i, image in enumerate(images, 1):
        print(f"正在处理第 {i}/{len(images)} 页...")

        # 添加空白幻灯片
        blank_slide_layout = prs.slide_layouts[6]  # 6是完全空白的布局
        slide = prs.slides.add_slide(blank_slide_layout)

        # 将PIL图片转换为字节流
        img_byte_arr = BytesIO()
        image.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)

        # 将图片添加到幻灯片，占满整个页面
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height

        slide.shapes.add_picture(img_byte_arr, left, top, width=width, height=height)

    # 保存PPT
    prs.save(output_ppt_path)
    print(f"\n✓ PPT已成功保存到: {output_ppt_path}")
    print(f"  共 {len(images)} 页")

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description='将PDF的每一页转换为高清图片并插入到PPT中',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  %(prog)s input.pdf output.pptx
  %(prog)s out/index.pdf out/presentation.pptx
        """
    )
    parser.add_argument('pdf_file', help='输入的PDF文件路径')
    parser.add_argument('output_file', nargs='?', help='输出的PPT文件路径（可选，默认为同名.pptx）')
    parser.add_argument('--dpi', type=int, default=200, help='图片DPI（默认200，越高越清晰）')

    args = parser.parse_args()

    pdf_file = args.pdf_file
    output_file = args.output_file

    # 如果没有指定输出文件，使用同名.pptx
    if not output_file:
        base = os.path.splitext(pdf_file)[0]
        output_file = f"{base}.pptx"

    # 检查PDF是否存在
    if not os.path.exists(pdf_file):
        print(f"错误: 找不到PDF文件 {pdf_file}")
        exit(1)

    # 确保输出目录存在
    output_dir = os.path.dirname(output_file)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    # 执行转换
    pdf_to_ppt(pdf_file, output_file)
