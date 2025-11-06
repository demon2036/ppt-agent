#!/usr/bin/env python3
"""检查PPT中的图片"""

from pptx import Presentation
import os

pptx_file = "WebLeaper专业组会展示.pptx"
prs = Presentation(pptx_file)

print(f"检查PPT: {pptx_file}")
print("=" * 60)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n幻灯片 {i}:")
    has_image = False

    for shape in slide.shapes:
        print(f"  - Shape类型: {shape.shape_type} ({type(shape).__name__})")

        if shape.shape_type == 13:  # Picture
            has_image = True
            try:
                image = shape.image
                print(f"    ✓ 找到图片! 格式: {image.ext}, 大小: {len(image.blob)} bytes")
            except Exception as e:
                print(f"    ✗ 无法读取图片: {e}")

        if hasattr(shape, "text") and shape.text.strip():
            text_preview = shape.text.strip()[:50]
            print(f"    文本: {text_preview}...")

    if not has_image:
        print(f"  ℹ️  此页无图片")

print("\n" + "=" * 60)
