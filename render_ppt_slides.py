#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
使用python-pptx读取PPT并使用Pillow渲染为图片
"""

from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import os
import textwrap

# 幻灯片尺寸 (16:9比例，高分辨率)
WIDTH = 1920
HEIGHT = 1080
BG_COLOR = (255, 255, 255)
TITLE_COLOR = (73, 31, 151)  # 紫色
TEXT_COLOR = (48, 58, 82)    # 深蓝色
SUBTITLE_COLOR = (100, 100, 100)

def wrap_text(text, font, max_width, draw):
    """文本换行处理"""
    words = text.split()
    lines = []
    current_line = []

    for word in words:
        test_line = ' '.join(current_line + [word])
        bbox = draw.textbbox((0, 0), test_line, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current_line.append(word)
        else:
            if current_line:
                lines.append(' '.join(current_line))
                current_line = [word]
            else:
                lines.append(word)

    if current_line:
        lines.append(' '.join(current_line))

    return lines

def render_slide(slide, slide_num, total_slides, output_path):
    """渲染单张幻灯片"""
    img = Image.new('RGB', (WIDTH, HEIGHT), BG_COLOR)
    draw = ImageDraw.Draw(img)

    # 尝试加载更好的字体，如果失败则使用默认字体
    try:
        title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 60)
        text_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
        small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28)
    except:
        title_font = None
        text_font = None
        small_font = None

    y_pos = 100

    # 处理标题
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()

            # 判断是否为标题（通常在上方且字体较大）
            if shape.top < HEIGHT // 3:  # 上方1/3区域认为是标题区
                # 绘制标题
                lines = wrap_text(text, title_font or text_font, WIDTH - 200, draw)
                for line in lines:
                    draw.text((100, y_pos), line, fill=TITLE_COLOR, font=title_font)
                    y_pos += 80
                y_pos += 40  # 标题后额外空间
            else:
                # 绘制内容文本
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        # 检查是否为项目符号
                        if line.strip().startswith('•') or line.strip().startswith('-'):
                            x_offset = 150
                        else:
                            x_offset = 100

                        wrapped_lines = wrap_text(line, text_font or small_font, WIDTH - 250, draw)
                        for wrapped_line in wrapped_lines:
                            if y_pos < HEIGHT - 150:  # 确保不超出边界
                                draw.text((x_offset, y_pos), wrapped_line, fill=TEXT_COLOR, font=text_font or small_font)
                                y_pos += 50

    # 绘制页码
    page_text = f"{slide_num} / {total_slides}"
    draw.text((WIDTH - 200, HEIGHT - 80), page_text, fill=SUBTITLE_COLOR, font=small_font)

    # 保存图片
    img.save(output_path, 'PNG', quality=95)

def main():
    """主函数"""
    pptx_file = "WebLeaper专业组会展示.pptx"
    output_dir = "screenshots"

    if not os.path.exists(pptx_file):
        print(f"错误: 找不到文件 {pptx_file}")
        return

    # 创建输出目录
    os.makedirs(output_dir, exist_ok=True)

    # 打开PPT
    print("正在读取PPT文件...")
    prs = Presentation(pptx_file)
    total_slides = len(prs.slides)

    print(f"找到 {total_slides} 张幻灯片")
    print("="*60)
    print("正在渲染幻灯片...")

    # 渲染每一张幻灯片
    for i, slide in enumerate(prs.slides, 1):
        output_path = os.path.join(output_dir, f"slide_{i:02d}.png")
        render_slide(slide, i, total_slides, output_path)
        print(f"  ✓ 已渲染第 {i}/{total_slides} 页: {output_path}")

    print("="*60)
    print(f"✓ 所有 {total_slides} 张幻灯片已成功渲染!")
    print(f"✓ 保存在: {output_dir}/")

if __name__ == "__main__":
    main()
