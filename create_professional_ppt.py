#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
WebLeaper 专业组会展示PPT生成脚本 - 包含论文图片和详细数据
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 颜色定义
TITLE_COLOR = RGBColor(73, 31, 151)  # 紫色
TEXT_COLOR = RGBColor(48, 58, 82)    # 深蓝色
HIGHLIGHT_COLOR = RGBColor(0, 102, 204)  # 蓝色
SUCCESS_COLOR = RGBColor(0, 153, 76)     # 绿色

def add_title_slide(prs):
    """标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "WebLeaper"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Empowering Efficiency and Efficacy in WebAgent\nvia Enabling Info-Rich Seeking"
    for p in subtitle_frame.paragraphs:
        p.font.size = Pt(22)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = TEXT_COLOR

    # 机构信息
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.6))
    author_frame = author_box.text_frame
    author_frame.text = "Tongyi Lab, Alibaba Group"
    p = author_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(100, 100, 100)

def add_slide_with_title(prs, title):
    """创建带标题的幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.LEFT

    return slide

def add_bullet_points(slide, points, top=1.5, left=0.8, width=8.5):
    """添加项目符号列表"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(18)
        p.space_before = Pt(8)
        p.font.color.rgb = TEXT_COLOR

def add_image_slide(prs, title, image_path, caption=""):
    """添加图片幻灯片"""
    slide = add_slide_with_title(prs, title)

    if os.path.exists(image_path):
        # 添加图片
        left = Inches(1.5)
        top = Inches(1.8)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(7))

        # 添加图注
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8.5), Inches(0.6))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            p = caption_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = RGBColor(100, 100, 100)
    else:
        # 如果图片不存在，显示占位符
        placeholder_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        placeholder_frame = placeholder_box.text_frame
        placeholder_frame.text = f"[图片: {os.path.basename(image_path)}]"
        p = placeholder_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(150, 150, 150)

def main():
    """主函数：生成完整的专业PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 图片路径前缀
    fig_path = "figures_png/"

    print("开始生成WebLeaper专业组会PPT...")
    print("=" * 70)

    # 1. 标题页
    print("  [1/20] 生成标题页...")
    add_title_slide(prs)

    # 2. 研究背景
    print("  [2/20] 生成研究背景页...")
    slide = add_slide_with_title(prs, "研究背景")
    add_bullet_points(slide, [
        "• LLM-based agents正在变革AI领域，推动复杂问题的解决",
        "• 信息搜索(Information Seeking, IS)是agent认知自主性的核心能力",
        "• 商业系统：OpenAI Deep Research, Google Gemini, Perplexity AI, Kimi-Researcher",
        "• 现有研究：主要关注搜索深度，较少关注搜索效率",
        "• 初步实验发现：当前LLM agents的搜索效率较低"
    ])

    # 3. 核心问题
    print("  [3/20] 生成核心问题页...")
    slide = add_slide_with_title(prs, "核心问题：低搜索效率")
    add_bullet_points(slide, [
        "• 问题现象：有效动作分布峰值约0.04，大部分动作无效",
        "• 低效行为表现：",
        "  - 冗余的查询重构",
        "  - 检索不相关的信息",
        "  - 不必要的长搜索链",
        "• 后果：增加计算和时间成本，限制整体IS性能",
        "• 根本原因：训练任务中目标实体过于稀疏"
    ])

    # 4. 有效动作分布图
    print("  [4/20] 生成有效动作分布图...")
    add_image_slide(prs, "有效动作分布",
                   fig_path + "infogain_cover_rate_kde.png",
                   "基于GPT模型在合成IS任务上的有效动作分布")

    # 5. 解决方案概览
    print("  [5/20] 生成解决方案概览...")
    slide = add_slide_with_title(prs, "解决方案：WebLeaper框架")
    add_bullet_points(slide, [
        "• 目标1：构建包含更多目标实体的新IS任务",
        "• 目标2：生成准确且高效的解决轨迹",
        "• 核心创新：",
        "  - 将IS建模为树状推理任务",
        "  - 提出三种数据集变体：Basic, Union, Reverse-Union",
        "  - 基于ISR和ISE指标过滤训练轨迹",
        "  - 设计混合奖励系统用于强化学习"
    ])

    # 6. 方法概览图
    print("  [6/20] 生成方法概览图...")
    add_image_slide(prs, "WebLeaper方法概览",
                   fig_path + "overview.png",
                   "(a) Basic: 单一信息源 | (b) Union: 多源融合 | (c) Reverse-Union: 反向推理")

    # 7. Version-I: Basic
    print("  [7/20] 生成Basic版本页...")
    slide = add_slide_with_title(prs, "数据合成 Version-I: Basic")
    add_bullet_points(slide, [
        "• 推理结构：树状结构 T_i，节点为实体，边为关系",
        "• 数据来源：Wikipedia表格（约200万个）",
        "• 三层结构：",
        "  - 根节点：表格标题中的实体（问题实体）",
        "  - 第二层：主键列的实体（如诺贝尔奖得主姓名）",
        "  - 第三层：其他列的属性（如国家、年份）",
        "• 特点：从单一信息源构建简单推理树"
    ])

    # 8. Version-II: Union
    print("  [8/20] 生成Union版本页...")
    slide = add_slide_with_title(prs, "数据合成 Version-II: Union")
    add_bullet_points(slide, [
        "• 动机：提高结构复杂度，跨越多个信息源",
        "• 方法：合并具有相似主题和结构的推理子树",
        "• 建模：最大二分团枚举(Maximal Biclique Enumeration)",
        "• 示例：",
        "  \"哪些作者同时获得了诺贝尔文学奖和布克奖？\"",
        "  - 需要识别两组获奖者（中间目标实体）",
        "  - 找到它们的交集（最终目标实体）",
        "• 特点：强制模型整合分散的互补证据"
    ])

    # 9. Version-III: Reverse-Union
    print("  [9/20] 生成Reverse-Union版本页...")
    slide = add_slide_with_title(prs, "数据合成 Version-III: Reverse-Union")
    add_bullet_points(slide, [
        "• 动机：防止简单的关键词搜索捷径",
        "• 反向推理流程：",
        "  1. Deductive Fuzz: 提供模糊线索（第三层实体）",
        "     \"1980年代获奖者，写了关于英国男孩困在荒岛的小说\"",
        "  2. 推断锚点实体：从线索推导出\"William Golding\"",
        "  3. Union-based Search: 使用锚点属性启动新搜索",
        "     \"找到所有同为英国人且同时获得两个奖项的作者\"",
        "• 特点：强制多步推理，提升规划和决策能力"
    ])

    # 10. 轨迹构造
    print("  [10/20] 生成轨迹构造页...")
    slide = add_slide_with_title(prs, "信息引导的轨迹构造")
    add_bullet_points(slide, [
        "• 工具配置：",
        "  - Search(queries, filter_year): Google搜索，支持时间过滤",
        "  - Visit(urls, goal): 访问URL并返回摘要",
        "• 过滤准则：",
        "  1. Coverage Criterion: ISR > α（覆盖率阈值）",
        "     ISR = |R ∩ O| / |R|（信息搜索率）",
        "  2. Efficiency Criterion: ISE > β（效率阈值）",
        "     ISE只计算Visit动作中的目标实体获取",
        "• 目标：保留准确且高效的轨迹作为训练数据"
    ])

    # 11. 强化学习
    print("  [11/20] 生成强化学习页...")
    slide = add_slide_with_title(prs, "混合奖励系统与强化学习")
    add_bullet_points(slide, [
        "• 挑战：",
        "  - 二元奖励过于稀疏（需要完美收集几十个实体）",
        "  - 精确匹配过于脆弱（\"USA\" vs \"United States\"）",
        "  - LLM-as-Judge成本高昂且不稳定",
        "• Hybrid Reward System：",
        "  - 实体密集型任务：基于F-score的细粒度奖励",
        "    R = (1+ω²) · (P·Rc) / (ω²P + Rc)",
        "  - 现有benchmark：保留原始奖励函数",
        "• 优化算法：GRPO (Group Relative Policy Optimization)"
    ])

    # 12. 实验设置
    print("  [12/20] 生成实验设置页...")
    slide = add_slide_with_title(prs, "实验设置")
    add_bullet_points(slide, [
        "• 评估基准（5个）：",
        "  - BrowseComp: 浏览和理解能力",
        "  - GAIA: 通用AI助手评估（103个文本子集）",
        "  - xbench-DeepSearch: 深度搜索能力",
        "  - Seal-0: 结构化信息抽取",
        "  - WideSearch: 广度搜索能力",
        "• 基线模型：",
        "  - 商用: Claude-4-Sonnet, OpenAI-o3, OpenAI DeepResearch",
        "  - 开源: WebSailor, WebDancer, Kimi-K2等10+模型",
        "• Base模型: Qwen3-30B-A3B-Thinking-2507"
    ])

    # 13. 主要结果图
    print("  [13/20] 生成主要结果图...")
    add_image_slide(prs, "主要实验结果 (Comprehensive Setting)",
                   fig_path + "main_results.png",
                   "在综合训练设置下的结果。所有分数均为3次运行的平均值")

    # 14. 主要结果表格
    print("  [14/20] 生成主要结果表格页...")
    slide = add_slide_with_title(prs, "详细实验结果对比")
    add_bullet_points(slide, [
        "• Base Setting关键结果：",
        "  WebLeaper-Reverse-Union (Base):",
        "  - BrowseComp: 23.0% | GAIA: 67.0%",
        "  - xbench-DS: 66.0% | Seal-0: 37.2%",
        "",
        "• Comprehensive Setting关键结果：",
        "  WebLeaper-Reverse-Union (Comprehensive):",
        "  - BrowseComp: 38.8% (最佳开源) | GAIA: 73.2%",
        "  - xbench-DS: 72.0% | WideSearch SR: 4.0%",
        "",
        "• 超越多个商用系统，成为最佳开源方案"
    ])

    # 15. 消融实验
    print("  [15/20] 生成消融实验页...")
    slide = add_slide_with_title(prs, "消融实验：不同数据源的影响")
    add_bullet_points(slide, [
        "• 对比baseline: WebSailor-V2 (5k/10k样本)",
        "",
        "• Basic-5k: 平均-7.64",
        "  - 任务过于简单，模型学会捷径而非真正的信息整合",
        "",
        "• Union-5k: 平均+3.26 ✓",
        "  - 多源整合提高复杂度，强制分散证据推理",
        "",
        "• Reverse-Union-10k: 平均+4.34 ✓✓",
        "  - 反向推理增加规划难度，显著提升决策能力",
        "  - WideSearch: +10.92 | xbench-DS: +6.00"
    ])

    # 16. 轨迹构造消融
    print("  [16/20] 生成轨迹构造消融图...")
    add_image_slide(prs, "信息引导轨迹构造策略的消融实验",
                   fig_path + "traj_ablation_bar_plot.png",
                   "对比ISR-Only, ISE-Only, 和ISR+ISE的效果")

    # 17. 效率与效果联合提升
    print("  [17/20] 生成效率与效果图...")
    add_image_slide(prs, "效率与效果的联合提升",
                   fig_path + "tool_call_performance_scatter.png",
                   "WebLeaper在效果和效率上均优于WebSailor-V2基线")

    # 18. RL结果
    print("  [18/20] 生成RL结果页...")
    slide = add_slide_with_title(prs, "强化学习结果")
    add_bullet_points(slide, [
        "• SFT+RL vs SFT-only (Comprehensive Setting):",
        "",
        "  BrowseComp: 37.8 → 38.8 (+1.0)",
        "  GAIA:       69.9 → 73.2 (+3.3) ✓",
        "  xbench-DS:  69.0 → 72.0 (+3.0) ✓",
        "  WideSearch SR: 1.5 → 4.0 (+2.5) ✓✓",
        "  WideSearch Row F1: 23.0 → 31.0 (+8.0) ✓✓",
        "",
        "• Reward曲线显示稳定持续上升趋势",
        "• 在135步时因web访问资源耗尽而终止，仍有潜力"
    ])

    # 19. RL奖励曲线
    print("  [19/20] 生成RL奖励曲线图...")
    add_image_slide(prs, "强化学习训练曲线",
                   fig_path + "reward_curve_compact.png",
                   "混合奖励系统的训练曲线，显示稳定的奖励增长")

    # 20. 主要贡献
    print("  [20/20] 生成主要贡献页...")
    slide = add_slide_with_title(prs, "主要贡献与创新")
    add_bullet_points(slide, [
        "• 创新性地将IS建模为树状推理问题",
        "  - 在有限上下文中嵌入更多目标实体",
        "",
        "• 设计三种任务合成变体（Basic, Union, Reverse-Union）",
        "  - 系统性提升训练数据的质量和多样性",
        "",
        "• 提出ISR和ISE指标及轨迹过滤策略",
        "  - 同时优化准确性和效率",
        "",
        "• 设计混合奖励系统用于强化学习",
        "  - 解决实体密集型任务的奖励稀疏问题",
        "",
        "• 在5个benchmark上验证有效性和泛化能力"
    ])

    # 21. 数据统计
    print("  [21/23] 生成数据统计页...")
    slide = add_slide_with_title(prs, "数据集统计")
    add_bullet_points(slide, [
        "• Wikipedia数据：",
        "  - 爬取约200万个表格",
        "  - 多阶段清洗，保留大型、结构良好的表格",
        "",
        "• 合成任务数量：",
        "  - Basic: 5,000样本",
        "  - Union: 5,000样本",
        "  - Reverse-Union: 10,000样本",
        "",
        "• 训练配置：",
        "  - 与5,000个WebSailor-V2数据混合",
        "  - ISR阈值α = 0.3, ISE阈值β = 0.1"
    ])

    # 22. 总结
    print("  [22/23] 生成总结页...")
    slide = add_slide_with_title(prs, "总结")
    add_bullet_points(slide, [
        "• WebLeaper成功解决了IS agents搜索效率低的问题",
        "",
        "• 树状推理建模实现了更高密度的实体嵌入",
        "",
        "• 三种任务变体系统性地提升了任务复杂度和多样性",
        "",
        "• 信息引导的轨迹构造确保了准确性和效率的双重优化",
        "",
        "• 混合奖励系统有效解决了RL中的奖励稀疏问题",
        "",
        "• 实验验证了在效率和效果上的显著提升",
        "",
        "• 成为当前最佳的开源信息搜索agent"
    ])

    # 23. Q&A页
    print("  [23/23] 生成Q&A页...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    thanks_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "谢谢！\n\nQ & A"
    for i, p in enumerate(thanks_frame.paragraphs):
        p.font.size = Pt(60 if i == 0 else 50)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = TITLE_COLOR

    # 保存PPT
    output_file = 'WebLeaper专业组会展示.pptx'
    prs.save(output_file)
    print("=" * 70)
    print(f"✓ PPT生成成功: {output_file}")
    print(f"✓ 总共生成 {len(prs.slides)} 张幻灯片")
    print("=" * 70)

if __name__ == "__main__":
    main()
