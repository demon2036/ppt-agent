#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
增强的PPT渲染脚本 - 包含图片渲染
"""

from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import os
import io

# 幻灯片尺寸 (16:9比例，高分辨率)
WIDTH = 1920
HEIGHT = 1080
BG_COLOR = (255, 255, 255)
TITLE_COLOR = (73, 31, 151)  # 紫色
TEXT_COLOR = (48, 58, 82)    # 深蓝色

def render_slide(slide, slide_num, total_slides, output_path):
    """渲染单张幻灯片，包括图片"""
    img = Image.new('RGB', (WIDTH, HEIGHT), BG_COLOR)
    draw = ImageDraw.Draw(img)

    # 加载字体
    try:
        title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 60)
        text_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 30)
        small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
    except:
        title_font = text_font = small_font = None

    # 首先渲染图片
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture type
            try:
                # 获取图片
                image = shape.image
                image_bytes = image.blob

                # 从字节加载图片
                pil_image = Image.open(io.BytesIO(image_bytes))

                # 计算位置和大小（从EMU转换为像素）
                left = int(shape.left * WIDTH / slide.slide_layout.width)
                top = int(shape.top * HEIGHT / slide.slide_layout.height)
                width = int(shape.width * WIDTH / slide.slide_layout.width)
                height = int(shape.height * HEIGHT / slide.slide_layout.height)

                # 调整图片大小
                pil_image = pil_image.resize((width, height), Image.Resampling.LANCZOS)

                # 粘贴图片到画布
                img.paste(pil_image, (left, top))

            except Exception as e:
                # 如果图片处理失败，绘制一个占位框
                try:
                    left = int(shape.left * WIDTH / slide.slide_layout.width)
                    top = int(shape.top * HEIGHT / slide.slide_layout.height)
                    width = int(shape.width * WIDTH / slide.slide_layout.width)
                    height = int(shape.height * HEIGHT / slide.slide_layout.height)
                    draw.rectangle([left, top, left + width, top + height],
                                 outline=(200, 200, 200), width=2)
                    draw.text((left + 10, top + 10), "[Image]",
                            fill=(150, 150, 150), font=small_font)
                except:
                    pass

    # 然后渲染文本
    y_pos = 80
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()

            # 跳过图片的alt text
            if shape.shape_type == 13:
                continue

            try:
                # 获取形状位置
                shape_top = int(shape.top * HEIGHT / slide.slide_layout.height)
                shape_left = int(shape.left * WIDTH / slide.slide_layout.width)
                shape_width = int(shape.width * WIDTH / slide.slide_layout.width)

                # 判断是否为标题（通常在上方且字体较大）
                if shape_top < HEIGHT // 3:
                    # 标题
                    lines = text.split('\n')
                    for line in lines[:2]:  # 只取前两行
                        if line.strip():
                            draw.text((shape_left, shape_top), line.strip(),
                                    fill=TITLE_COLOR, font=title_font)
                            shape_top += 80
                else:
                    # 内容文本
                    lines = text.split('\n')
                    for line in lines:
                        if line.strip() and shape_top < HEIGHT - 100:
                            draw.text((shape_left, shape_top), line.strip(),
                                    fill=TEXT_COLOR, font=text_font)
                            shape_top += 45
            except:
                pass

    # 绘制页码
    draw.text((WIDTH - 200, HEIGHT - 80), f"{slide_num} / {total_slides}",
             fill=(150, 150, 150), font=small_font)

    # 保存图片
    img.save(output_path, 'PNG', quality=95)

def main():
    """主函数"""
    pptx_file = "WebLeaper专业组会展示.pptx"
    output_dir = "screenshots"

    if not os.path.exists(pptx_file):
        print(f"错误: 找不到文件 {pptx_file}")
        return

    # 创建输出目录
    os.makedirs(output_dir, exist_ok=True)

    # 打开PPT
    print("正在读取PPT文件...")
    prs = Presentation(pptx_file)
    total_slides = len(prs.slides)

    print(f"找到 {total_slides} 张幻灯片")
    print("="*60)
    print("正在渲染幻灯片（包含图片）...")

    # 渲染每一张幻灯片
    for i, slide in enumerate(prs.slides, 1):
        output_path = os.path.join(output_dir, f"slide_{i:02d}.png")
        render_slide(slide, i, total_slides, output_path)
        print(f"  ✓ 已渲染第 {i}/{total_slides} 页: {output_path}")

    print("="*60)
    print(f"✓ 所有 {total_slides} 张幻灯片已成功渲染!")
    print(f"✓ 保存在: {output_dir}/")

if __name__ == "__main__":
    main()
