#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
WebLeaper 组会展示PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加标题
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "WebLeaper"

    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(73, 31, 151)  # 紫色
    p.alignment = PP_ALIGN.CENTER

    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Empowering Efficiency and Efficacy in WebAgent\nvia Enabling Info-Rich Seeking"

    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.color.rgb = RGBColor(48, 58, 82)

    # 添加作者信息
    author_box = slide.shapes.add_textbox(left, top + Inches(2.5), width, Inches(0.8))
    author_frame = author_box.text_frame
    author_frame.text = "Tongyi Lab, Alibaba Group"

    p = author_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(100, 100, 100)

def add_content_slide(prs, title, content_points):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局

    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(73, 31, 151)

    # 添加内容
    left = Inches(0.8)
    top = Inches(2.0)
    width = Inches(8.5)
    height = Inches(4.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(content_points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.space_before = Pt(12)
        p.level = 0

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """添加两栏内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(73, 31, 151)
    p.alignment = PP_ALIGN.CENTER

    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(4.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    # 左栏标题
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_after = Pt(12)

    # 左栏内容
    for point in left_points:
        p = left_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.space_before = Pt(8)

    # 右栏
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(4.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    # 右栏标题
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 76)
    p.space_after = Pt(12)

    # 右栏内容
    for point in right_points:
        p = right_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.space_before = Pt(8)

def main():
    """主函数：生成完整的PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 标题页
    add_title_slide(prs)

    # 2. 研究背景与动机
    add_content_slide(prs, "研究背景与动机", [
        "LLM-based agents正在变革开放式问题解决方法",
        "信息搜索(Information Seeking)是自主推理和决策的核心能力",
        "现有问题：当前IS agents存在搜索效率低的问题",
        "根本原因：训练任务中目标实体稀疏，限制了代理学习高效搜索行为的机会"
    ])

    # 3. 核心问题
    add_content_slide(prs, "核心问题", [
        "问题1：如何提高信息搜索代理的搜索效率？",
        "问题2：如何在有限的上下文中嵌入更多的目标实体？",
        "问题3：如何生成既准确又高效的解决方案轨迹？",
        "挑战：需要同时优化正确性和搜索性能"
    ])

    # 4. 解决方案：WebLeaper框架
    add_content_slide(prs, "解决方案：WebLeaper框架", [
        "将信息搜索形式化为树状推理问题",
        "使用Wikipedia表格构建高覆盖率的IS任务",
        "提出三种任务合成变体：Basic、Union、Reverse-Union",
        "仅保留准确且高效的训练轨迹",
        "系统性地提高IS效率和效果"
    ])

    # 5. 方法详解
    add_two_column_slide(prs, "方法详解",
        "任务构建", [
            "Basic：基础实体搜索任务",
            "Union：联合多实体搜索",
            "Reverse-Union：逆向联合搜索",
            "利用Wikipedia结构化数据"
        ],
        "轨迹优化", [
            "过滤低效搜索路径",
            "保留准确答案",
            "优化搜索步骤数",
            "平衡效率与效果"
        ]
    )

    # 6. 实验结果
    add_content_slide(prs, "实验结果", [
        "在5个IS基准上进行了广泛实验：",
        "  - BrowserComp (英文/中文)",
        "  - GAIA",
        "  - xbench-DeepSearch",
        "  - WideSearch",
        "  - Seal-0",
        "在基础和综合设置下，均实现了效率和效果的双重提升",
        "相比强基线模型表现出持续的改进"
    ])

    # 7. 主要贡献
    add_content_slide(prs, "主要贡献", [
        "✓ 提出WebLeaper框架，系统性解决IS效率问题",
        "✓ 创新性地将IS建模为树状推理问题",
        "✓ 设计了三种任务合成变体，提高训练数据质量",
        "✓ 提出轨迹过滤策略，同时优化准确性和效率",
        "✓ 在多个benchmark上验证了方法的有效性"
    ])

    # 8. 技术亮点
    add_two_column_slide(prs, "技术亮点",
        "创新点", [
            "树状推理建模",
            "高覆盖率任务构建",
            "多变体任务合成",
            "双重优化策略"
        ],
        "优势", [
            "提高搜索效率",
            "增强泛化能力",
            "保证答案准确性",
            "降低计算成本"
        ]
    )

    # 9. 实验设置
    add_content_slide(prs, "实验设置", [
        "数据来源：Wikipedia结构化表格",
        "评估基准：5个不同的IS benchmarks",
        "评估指标：准确率、搜索效率、成功率",
        "对比基线：多个SOTA的深度研究模型/系统",
        "实验环境：Basic设置和Comprehensive设置"
    ])

    # 10. 局限性与未来工作
    add_two_column_slide(prs, "局限性与未来工作",
        "当前局限", [
            "依赖Wikipedia数据源",
            "任务类型相对受限",
            "需要人工设计变体"
        ],
        "未来方向", [
            "扩展到更多数据源",
            "自动化任务生成",
            "多模态信息搜索",
            "实时搜索优化"
        ]
    )

    # 11. 总结
    add_content_slide(prs, "总结", [
        "WebLeaper通过创新的框架设计，成功解决了IS agents效率低的问题",
        "树状推理建模使得在有限上下文中嵌入更多目标实体成为可能",
        "三种任务合成变体系统性地提升了训练数据的质量和多样性",
        "轨迹过滤策略确保了准确性和效率的双重优化",
        "实验结果证明了方法在多个benchmark上的有效性和泛化能力"
    ])

    # 12. 致谢页
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    thanks_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thanks!\n\nQ & A"

    for i, paragraph in enumerate(thanks_frame.paragraphs):
        paragraph.font.size = Pt(60) if i == 0 else Pt(48)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.color.rgb = RGBColor(73, 31, 151)

    # 保存PPT
    prs.save('WebLeaper组会展示.pptx')
    print("✓ PPT生成成功: WebLeaper组会展示.pptx")
    print(f"✓ 总共生成 {len(prs.slides)} 张幻灯片")

if __name__ == "__main__":
    main()
